#!/usr/bin/env python3
"""Remove figures that appear twice in the deck, keeping the better placement.

A figure shown on two slides costs width on both and buys nothing on the second.
Dropping the weaker placement is the only lever that helps the slides where the
figures' combined need genuinely exceeds the room available — no rendering choice
can fix a slide whose figures need 17.5 in of width in a 12.9 in column.

Each decision below keeps the instance where the figure carries the argument and
drops the one where it is a restatement. Nothing is removed from the deck as a
whole: every figure listed still appears on its kept slide.

Slide 1 is not a duplicate but a MISMATCH: its caption reads "Distributed-
temperature waterfall T(x,t)" while the picture is the six-panel space-time
solution. It is corrected to the figure the caption describes, which also stops
19_spacetime_fields appearing three times.
"""
import hashlib
import os
import sys

from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
CASE = os.path.abspath(os.path.join(HERE, ".."))
DECK = "/mnt/c/Users/user/Desktop/slides3.pptx"

#  (slide number, figure to remove, why)
DROP = [
    (24, "24_temperature_gradient.png",
     "kept on 25; slide 24 carries three figures needing 17.5 in of a 12.9 in row"),
    (33, "06_deposit.png",
     "kept on 24, where the deposit is the subject rather than a comparison"),
    (28, "12_mitigation_comparison.png",
     "kept on 33, the cross-scenario slide the comparison belongs to"),
    (14, "compo_pvt.png",
     "kept on 27, where it is the only figure and reads at full width"),
    (37, "27_wellposedness_map.png",
     "kept on 29, where it is the only figure"),
    (35, "friction_validation.png",
     "kept on 12; leaves hydrate_validation the whole column on 35"),
]

#  slide -> (figure now shown, figure the caption actually describes)
RETARGET = [(1, "19_spacetime_fields.png", "23_dts_thermal_waterfall.png")]


def index():
    idx = {}
    for d in sorted(os.listdir(CASE)):
        p = os.path.join(CASE, d)
        if not (d.startswith("outputs") and os.path.isdir(p)):
            continue
        for f in os.listdir(p):
            if f.endswith((".png", ".gif")):
                try:
                    h = hashlib.sha256(open(os.path.join(p, f), "rb").read()).hexdigest()
                except OSError:
                    continue
                idx.setdefault(h, (f, os.path.join(p, f)))
    return idx


def find(case_dirs, name):
    """Any rendering of `name`, preferring a slide-legible one."""
    for d in sorted(os.listdir(CASE), key=lambda x: (not x.startswith("outputs_slides"), x)):
        p = os.path.join(CASE, d, name)
        if os.path.isfile(p):
            return p
    return None


def main(argv):
    deck = argv[0] if argv else DECK
    prs = Presentation(deck)
    idx = index()
    slides = list(prs.slides)
    removed = retargeted = 0

    for num, name, why in DROP:
        if num > len(slides):
            print(f"  [skip] slide {num} does not exist")
            continue
        s = slides[num - 1]
        for sh in list(s.shapes):
            if sh.__class__.__name__ != "Picture":
                continue
            got = idx.get(hashlib.sha256(sh.image.blob).hexdigest(), ("", ""))[0]
            if got == name:
                #  A figure does not sit alone: it has a white card behind it and a
                #  caption beneath. Leaving those behind turns the drop into an empty
                #  panel that still BLOCKS its neighbour from expanding — which is
                #  exactly the space the drop was meant to release.
                bx = (Emu(sh.left).inches, Emu(sh.top).inches,
                      Emu(sh.left).inches + Emu(sh.width).inches,
                      Emu(sh.top).inches + Emu(sh.height).inches)
                for o in list(s.shapes):
                    if o is sh or o.__class__.__name__ == "Picture":
                        continue
                    try:
                        a0, b0 = Emu(o.left).inches, Emu(o.top).inches
                        a1, b1 = a0 + Emu(o.width).inches, b0 + Emu(o.height).inches
                    except Exception:
                        continue
                    ov = (max(0.0, min(bx[2], a1) - max(bx[0], a0))
                          * max(0.0, min(bx[3], b1) - max(bx[1], b0)))
                    area = max((a1 - a0) * (b1 - b0), 1e-9)
                    txt = (o.text_frame.text.strip() if o.has_text_frame else "")
                    #  the card behind it: no text, and mostly covered by the figure
                    if not txt and ov / area > 0.22:
                        o._element.getparent().remove(o._element)
                    #  its caption: a short line directly beneath, left-aligned with it
                    elif (txt and len(txt) < 220 and 0 <= b0 - bx[3] < 0.7
                          and abs(a0 - bx[0]) < 1.6):
                        o._element.getparent().remove(o._element)
                sh._element.getparent().remove(sh._element)
                removed += 1
                print(f"  slide {num:>2}  dropped {name:<32} {why}")
                break
        else:
            print(f"  slide {num:>2}  {name} not found (already dropped?)")

    for num, old, new in RETARGET:
        s = slides[num - 1]
        src = find(CASE, new)
        if src is None:
            print(f"  slide {num:>2}  {new} not on disk; left as is")
            continue
        for sh in list(s.shapes):
            if sh.__class__.__name__ != "Picture":
                continue
            got = idx.get(hashlib.sha256(sh.image.blob).hexdigest(), ("", ""))[0]
            if got == old:
                L, T, W, H = sh.left, sh.top, sh.width, sh.height
                pic = s.shapes.add_picture(src, L, T, W, H)
                sh._element.addprevious(pic._element)
                sh._element.getparent().remove(sh._element)
                retargeted += 1
                print(f"  slide {num:>2}  {old} -> {new}  (matches the caption)")
                break

    #  A card can outlive its figure by more than the drop pass can see, because the
    #  refit moves the surviving pictures afterwards. Sweep the affected slides once
    #  more for figure-sized boxes that now hold neither text nor picture.
    orphans = 0
    for num in {n for n, _, _ in DROP}:
        s = slides[num - 1]
        pboxes = []
        for o in s.shapes:
            if o.__class__.__name__ == "Picture":
                a0, b0 = Emu(o.left).inches, Emu(o.top).inches
                pboxes.append((a0, b0, a0 + Emu(o.width).inches, b0 + Emu(o.height).inches))
        for o in list(s.shapes):
            if o.__class__.__name__ == "Picture":
                continue
            if o.has_text_frame and o.text_frame.text.strip():
                continue
            try:
                a0, b0 = Emu(o.left).inches, Emu(o.top).inches
                a1, b1 = a0 + Emu(o.width).inches, b0 + Emu(o.height).inches
            except Exception:
                continue
            w, h = a1 - a0, b1 - b0
            #  a figure card, not a rule or a divider or a full-width banner
            if not (1.5 <= w <= 9.0 and 0.8 <= h <= 5.0):
                continue
            #  "does a figure actually sit in this card" — the same question
            #  check_slides asks. An area-overlap test disagreed with it and left a
            #  card standing that the checker still reported as an empty panel.
            if any(a0 - 0.05 <= (q[0] + q[2]) / 2 <= a1 + 0.05
                   and b0 - 0.05 <= (q[1] + q[3]) / 2 <= b1 + 0.05
                   for q in pboxes):
                continue
            o._element.getparent().remove(o._element)
            orphans += 1
            print(f"  slide {num:>2}  removed an orphaned {w:.2f}x{h:.2f} in card")

    prs.save(deck)
    print(f"\n  {removed} duplicate(s) dropped, {retargeted} retargeted, "
          f"{orphans} orphaned card(s) swept")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
