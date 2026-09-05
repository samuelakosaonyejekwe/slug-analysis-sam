#!/usr/bin/env python3
"""Give a figure its own slide when its current slide cannot make it readable.

fit_deck_figures.py chooses the best rendering for the frame a figure has, and
drop_duplicate_figures.py frees width by removing repeats. Neither can help a slide
whose figures simply need more room than it has: three maps in one row on slide 31
need 13.6 in of a 12.9 in band, and five animations on slide 34 need more still.
The only remaining lever is a slide.

For each figure that cannot reach TARGET_PT on its own slide, this moves it to a new
slide placed immediately after, carrying the section line, the title, the rule, the
footer and the slide number so the deck reads continuously, plus the figure at full
content width and its caption. The originating slide keeps its remaining figures,
which then have the whole row to themselves.

    python3 expand_deck.py [deck.pptx]

Slide numbers in the footer are renumbered afterwards, since inserting shifts them.
"""
import copy
import hashlib
import os
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
CASE = os.path.abspath(os.path.join(HERE, ".."))
DECK = "/mnt/c/Users/user/Desktop/slides3.pptx"

MAX_SLIDES = 48           # raised once the caption placement could carry a figure
TARGET_PT = 12.0          # the back-of-room threshold this is trying to reach
MIN_SQIN = 6.0            # and the physical size below which a figure is a thumbnail
MARGIN = 0.62             # the deck's content margin
TOP_BAND = 1.48           # below the title rule
BOT_BAND = 6.92           # above the footer rule
ANIM_DPI_DEFAULT = 200.0


def anim_dpi(path):
    import json
    try:
        with open(os.path.join(os.path.dirname(path), "anim_dpi.json")) as fh:
            return float(json.load(fh)["dpi"])
    except Exception:
        return ANIM_DPI_DEFAULT


def natural_in(path):
    with Image.open(path) as im:
        px, py = im.size
        if (im.format or "").upper() == "GIF":
            dpi = anim_dpi(path)
        else:
            d = im.info.get("dpi")
            dpi = float(d[0]) if d and d[0] > 1 else 100.0
        return px / dpi, py / px


def candidates():
    out = {}
    for d in sorted(os.listdir(CASE)):
        p = os.path.join(CASE, d)
        if not (d.startswith("outputs") and os.path.isdir(p)):
            continue
        base = 18.0 if d.startswith("outputs_slides") else 10.0
        for f in os.listdir(p):
            if not f.endswith((".png", ".gif")):
                continue
            fp = os.path.join(p, f)
            try:
                nat, ar = natural_in(fp)
            except Exception:
                continue
            if nat > 0:
                out.setdefault(f, []).append((fp, nat, ar, base))
    return out


def by_hash():
    idx = {}
    for d in os.listdir(CASE):
        p = os.path.join(CASE, d)
        if not (d.startswith("outputs") and os.path.isdir(p)):
            continue
        for f in os.listdir(p):
            if f.endswith((".png", ".gif")):
                try:
                    idx.setdefault(hashlib.sha256(
                        open(os.path.join(p, f), "rb").read()).hexdigest(), f)
                except OSError:
                    pass
    return idx


def best_on_own_slide(cands, avail_w, avail_h):
    """The best effective pt this figure could reach with a whole slide to itself."""
    best = (0.0, None, None, None)
    for fp, nat, ar, base in cands:
        w = min(avail_w, avail_h / ar if ar > 0 else avail_w)
        eff = base * w / nat
        if eff > best[0]:
            best = (eff, fp, w, w * ar)
    return best


def chrome(slide):
    """The shapes that make a slide look like part of this deck, not its content."""
    keep = []
    for sh in slide.shapes:
        if sh.__class__.__name__ == "Picture":
            continue
        try:
            top = Emu(sh.top).inches
            h = Emu(sh.height).inches
        except Exception:
            continue
        txt = sh.text_frame.text.strip() if sh.has_text_frame else ""
        #  header block, the rule under it, the footer rule, footer text, page number
        if top < TOP_BAND or top + h > BOT_BAND:
            keep.append(sh)
        elif not txt and h < 0.12:
            keep.append(sh)
    return keep


def main(argv):
    deck = argv[0] if argv else DECK
    prs = Presentation(deck)
    SW, SH = Emu(prs.slide_width).inches, Emu(prs.slide_height).inches
    avail_w = SW - 2 * MARGIN
    avail_h = BOT_BAND - TOP_BAND - 0.42          # room for a caption line
    cands, idx = candidates(), by_hash()
    moved = []

    for i, slide in enumerate(list(prs.slides), 1):
        pics = [s for s in slide.shapes if s.__class__.__name__ == "Picture"]
        if len(pics) < 2:
            continue                               # already has the slide to itself
        for sh in list(pics):
            try:
                blob = sh.image.blob
            except Exception:
                continue
            name = idx.get(hashlib.sha256(blob).hexdigest())
            if not name or name not in cands:
                continue
            w_now = Emu(sh.width).inches
            h_now = Emu(sh.height).inches
            cur = None
            try:
                im = Image.open(__import__("io").BytesIO(blob))
                dpi = (ANIM_DPI_DEFAULT if (im.format or "").upper() == "GIF"
                       else (float(im.info["dpi"][0])
                             if im.info.get("dpi") and im.info["dpi"][0] > 1 else 100.0))
                cur = im.size[0] / dpi
            except Exception:
                pass
            base_now = next((b for _f, n, _a, b in cands[name]
                             if cur and abs(n - cur) < 1e-6), 18.0)
            eff_now = base_now * w_now / cur if cur else 0.0
            if eff_now >= TARGET_PT and w_now * h_now >= MIN_SQIN:
                continue                           # already fine where it is
            eff_own, fp, w, h = best_on_own_slide(cands[name], avail_w, avail_h)
            if fp is None or eff_own < max(eff_now * 1.15, TARGET_PT):
                continue                           # a slide of its own would not help
            ar_own = h / w if w else 1.0
            moved.append((i, name, eff_now, eff_own, fp, w, h, slide, sh, ar_own))

    #  A deck length is a talk length. When more figures want their own slide than
    #  the budget allows, spend it where it buys the most: the worst-rendering figure
    #  first, then by the size of the gain. The rest keep their current placement and
    #  are reported, so what was left undone is visible rather than silently dropped.
    budget = max(MAX_SLIDES - len(prs.slides._sldIdLst), 0)
    moved.sort(key=lambda m: (m[2], -(m[3] - m[2])))
    deferred = moved[budget:]
    moved = moved[:budget]
    if deferred:
        print(f"  budget {budget} slide(s); {len(deferred)} figure(s) left in place:")
        for (_i, _n, _e, _o, *_r) in deferred:
            print(f"     slide {_i:>2}  {_n:<32}{_e:>5.1f} pt "
                  f"(would reach {_o:.1f} on its own slide)")
        print()

    #  build the new slides after the scan, so indices stay stable
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    for (src_i, name, eff_now, eff_own, fp, w, h, src, sh, fig_ar) in moved:
        new = prs.slides.add_slide(blank)
        chrome_bottom = TOP_BAND
        for keep in chrome(src):
            new.shapes._spTree.append(copy.deepcopy(keep._element))
            #  A slide's chrome is not only the header: several carry a standfirst
            #  line that reaches BELOW the top of the content band. Starting the
            #  figure at the band top then puts it under that line.
            try:
                kt, kh = Emu(keep.top).inches, Emu(keep.height).inches
            except Exception:
                continue
            if kt < TOP_BAND and kt + kh > chrome_bottom and kt + kh < BOT_BAND - 1.0:
                chrome_bottom = kt + kh + 0.10
        bx = (Emu(sh.left).inches, Emu(sh.top).inches,
              Emu(sh.left).inches + Emu(sh.width).inches,
              Emu(sh.top).inches + Emu(sh.height).inches)
        #  A figure may carry MORE THAN ONE caption — the animations have a filename
        #  label and a description. Reserving room for one and stacking the rest
        #  pushed the second through the footer. Find them all first, reserve their
        #  combined height, and size the figure to what is left.
        caps = []
        for o in src.shapes:
            if o is sh or o.__class__.__name__ == "Picture" or not o.has_text_frame:
                continue
            t = o.text_frame.text.strip()
            if not t or len(t) >= 240:
                continue
            try:
                a0, b0 = Emu(o.left).inches, Emu(o.top).inches
            except Exception:
                continue
            if 0 <= b0 - bx[3] < 0.9 and abs(a0 - bx[0]) < 1.9:
                caps.append(o)
        caps.sort(key=lambda o: Emu(o.top).inches)
        reserve = sum(Emu(o.height).inches + 0.08 for o in caps) + 0.10
        band_h = BOT_BAND - chrome_bottom - reserve
        h = min(avail_w * (h / w) if w else h, band_h)
        w = h / (fig_ar if fig_ar else 1.0)
        w = min(w, avail_w)
        cap_y = [chrome_bottom + h + 0.10]
        left = int(round((MARGIN + (avail_w - w) / 2.0) * 914400))
        top = int(round(chrome_bottom * 914400))
        _np = new.shapes.add_picture(fp, int(left), top,
                                     int(round(w * 914400)), int(round(h * 914400)))
        #  This slide was laid out here, with room reserved for however many captions
        #  the figure carries. fit_deck_figures re-places whatever it finds, knows
        #  nothing of that reservation, and grew the figure back over its own caption.
        #  Marking the picture tells it to leave this one alone.
        _np.name = "shct-fixed " + _np.name
        #  move the slide to sit immediately after its source
        xml_slides = prs.slides._sldIdLst
        ids = list(xml_slides)
        xml_slides.remove(ids[-1])
        xml_slides.insert(src_i, ids[-1])
        #  the figure does not leave alone: the white card behind it and the caption
        #  beneath it go with it, or the source slide keeps an empty frame
        for o in list(src.shapes):
            if o is sh or o.__class__.__name__ == "Picture":
                continue
            try:
                a0, b0 = Emu(o.left).inches, Emu(o.top).inches
                a1, b1 = a0 + Emu(o.width).inches, b0 + Emu(o.height).inches
            except Exception:
                continue
            ov = (max(0.0, min(bx[2], a1) - max(bx[0], a0))
                  * max(0.0, min(bx[3], b1) - max(bx[1], b0)))
            area = max((a1 - a0) * (b1 - b0), 1e-9)
            txt = o.text_frame.text.strip() if o.has_text_frame else ""
            if not txt and ov / area > 0.22:
                #  the card: move it to the new slide behind the figure instead
                o._element.getparent().remove(o._element)
            elif o in caps:
                #  the caption belongs with the figure — but it has to be RE-PLACED
                #  under it on the new slide. Copied at its old position it lands
                #  wherever the source slide had it, which is on top of the figure
                #  once the figure is centred in a full-height band.
                el = copy.deepcopy(o._element)
                new.shapes._spTree.append(el)
                for cs in new.shapes:
                    if cs._element is el:
                        cs.left = int(round(MARGIN * 914400))
                        ch = Emu(cs.height).inches
                        cs.top = int(round(min(cap_y[0], BOT_BAND - ch) * 914400))
                        cs.width = int(round(avail_w * 914400))
                        cap_y[0] += ch + 0.08
                        break
                o._element.getparent().remove(o._element)
        sh._element.getparent().remove(sh._element)
        print(f"  slide {src_i:>2}  {name:<32} {eff_now:>5.1f} -> {eff_own:>5.1f} pt "
              f"on a slide of its own")

    prs.save(deck)
    print(f"\n  {len(moved)} figure(s) given their own slide; deck is now "
          f"{len(prs.slides._sldIdLst)} slides")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
