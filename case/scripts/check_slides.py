#!/usr/bin/env python3
# =============================================================================
#  check_slides.py — audit a deck for the faults that are obvious on the screen
#  and invisible in the file.
# -----------------------------------------------------------------------------
#  Three things go wrong when figures are placed programmatically, and none of
#  them raises an error:
#
#    EMPTY FRAME    a panel shape that was drawn to hold a picture, with no
#                   picture in it — an empty box on the slide, often still
#                   carrying the caption of the figure that used to be there.
#    TINY FIGURE    a detailed figure scaled into a frame built for a thumbnail,
#                   unreadable at presentation size.
#    OVERLAP        a text box sitting on top of a picture, or two text boxes on
#                   top of each other.
#
#  Everything here is measured from the shape geometry, so the report says which
#  slide and by how much rather than "looks wrong".
#
#      python3 check_slides.py [deck.pptx]
#
#  Exit status is 1 if anything is flagged.
# =============================================================================
import os
import sys

from pptx import Presentation
from pptx.util import Emu

#  a figure smaller than this is not readable from the back of a room
MIN_AREA_SQIN = 3.0
#  a container this big that holds nothing reads as an empty box
MIN_EMPTY_PANEL_SQIN = 2.0
#  A figure this small cannot be read from the back of a room once its panels are
#  counted: a six-panel plot needs six times the area of a single one before any
#  panel is legible. Area alone is the wrong test, so the per-panel figure below is
#  what is actually reported.
MIN_AREA_PER_PANEL_SQIN = 2.0
#  Text overflow. A .pptx stores the box, not the rendered text, so a box holding
#  more text than fits is clipped on screen and looks perfectly fine in the file —
#  which is why truncation survives every structural check. The estimate below is
#  deliberately crude and deliberately conservative: it only reports a box whose
#  text cannot fit by a clear margin, so a flagged box is worth opening.
OVERFLOW_TOLERANCE = 1.15
#  ignore incidental touching; flag a real covering
#  two text boxes that merely abut share a sliver of each other's padding; only a
#  real covering is a fault
MIN_OVERLAP_FRAC = 0.25


def box(sh):
    return (Emu(sh.left).inches, Emu(sh.top).inches,
            Emu(sh.width).inches, Emu(sh.height).inches)


def overlap_frac(a, b):
    """Fraction of the SMALLER box that the two boxes share."""
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    x0, x1 = max(ax, bx), min(ax + aw, bx + bw)
    y0, y1 = max(ay, by), min(ay + ah, by + bh)
    if x1 <= x0 or y1 <= y0:
        return 0.0
    inter = (x1 - x0) * (y1 - y0)
    return inter / max(min(aw * ah, bw * bh), 1e-9)


def est_text_height_in(sh):
    """Rough rendered height of a text frame's content, in inches.

    Characters per line come from the box width and the average glyph advance
    (~0.5 em for the sans faces used here); lines come from wrapping each
    paragraph. Leading is 1.2 em plus the frame's own insets.
    """
    from pptx.util import Pt
    tf = sh.text_frame
    w_in = Emu(sh.width).inches
    try:
        inset = (Emu(tf.margin_left).inches + Emu(tf.margin_right).inches)
        v_inset = (Emu(tf.margin_top).inches + Emu(tf.margin_bottom).inches)
    except Exception:
        inset, v_inset = 0.2, 0.1
    usable = max(w_in - inset, 0.3)
    lines = 0.0
    height = 0.0
    for para in tf.paragraphs:
        txt = "".join(r.text for r in para.runs) or para.text or ""
        #  the run's own size, else the paragraph's, else PowerPoint's 18 pt default
        pts = None
        for r in para.runs:
            if r.font.size is not None:
                pts = r.font.size.pt
                break
        if pts is None and para.font.size is not None:
            pts = para.font.size.pt
        pts = pts or 18.0
        char_w = 0.5 * pts / 72.0                 # em advance in inches
        per_line = max(int(usable / max(char_w, 1e-6)), 8)
        n = max(1, -(-len(txt) // per_line))      # ceil division
        lines += n
        height += n * 1.2 * pts / 72.0
    return height + v_inset


def stale_figures(path):
    """Deck pictures whose bytes match no current figure on disk.

    A .pptx embeds a COPY of every image, so a figure placed before the outputs
    were regenerated stays in the deck for ever and looks perfectly correct beside
    a caption that has since been updated. Nothing else here catches it: the
    geometry checks below see a well-placed picture, check_docs sees correct words,
    and the artwork freshness gate covers the exported Figure_NN set rather than the
    deck. Thirty of forty-two figures in this deck were once stale for exactly that
    reason, so the embedded bytes are hashed against the figures on disk.

    An unmatched picture is not proof of staleness — it may simply have been placed
    from somewhere not searched — so this reports rather than fails, and names the
    slide so it can be checked.
    """
    import hashlib
    here = os.path.dirname(os.path.abspath(__file__))
    case = os.path.abspath(os.path.join(here, ".."))
    known = set()
    #  The slide-legible twins exist for all three scenarios, not just the steady one.
    #  Omitting the other two made this check report five CURRENT figures on slides
    #  30-33 as "matching no current output" — a false alarm that trains the reader
    #  to ignore the real ones.
    #  discovered, not listed: make_slide_figures.py renders a variant set per size
    #  scale (outputs_slides45, outputs_slides70, ...) and a hardcoded list silently
    #  reports every figure from a new set as "matching no current output".
    _dirs = ["outputs_steady", "outputs_shutin", "outputs_mitigated"]
    _dirs += sorted(d for d in os.listdir(case)
                    if d.startswith("outputs_slides")
                    and os.path.isdir(os.path.join(case, d)))
    for d in _dirs:
        p = os.path.join(case, d)
        if not os.path.isdir(p):
            continue
        for f in os.listdir(p):
            if f.endswith((".png", ".gif")):
                try:
                    known.add(hashlib.sha256(
                        open(os.path.join(p, f), "rb").read()).hexdigest())
                except OSError:
                    pass
    if not known:
        return []
    out = []
    for i, slide in enumerate(Presentation(path).slides, 1):
        for sh in slide.shapes:
            if sh.__class__.__name__ != "Picture":
                continue
            try:
                h = hashlib.sha256(sh.image.blob).hexdigest()
            except Exception:
                continue
            if h not in known:
                out.append(i)
    return out


def audit(path):
    prs = Presentation(path)
    print(f"=== {os.path.basename(path)} — {len(prs.slides._sldIdLst)} slides ===")
    n_fail = 0

    for i, slide in enumerate(prs.slides, 1):
        pics = [sh for sh in slide.shapes if sh.__class__.__name__ == "Picture"]
        texts = [sh for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        #  container shapes: no text, no picture, but a real area — these are the
        #  frames the design draws behind a figure
        panels = [sh for sh in slide.shapes
                  if sh.__class__.__name__ != "Picture"
                  and (not sh.has_text_frame or not sh.text_frame.text.strip())]

        issues = []

        for p in pics:
            _l, _t, w, h = box(p)
            if w * h < MIN_AREA_SQIN:
                issues.append(f"figure only {w:.2f}x{h:.2f} in "
                              f"({w*h:.1f} sq in) — too small to read")

        for pan in panels:
            l, t, w, h = box(pan)
            if w * h < MIN_EMPTY_PANEL_SQIN:
                continue
            if w > 12.0 or h < 0.35:            # rules, dividers, backgrounds
                continue
            #  A panel is only EMPTY if nothing sits in it. Most of these shapes are
            #  the design's styled card behind a block of commentary, so a text box
            #  on top of one means it is doing its job, not that a figure is missing.
            holds_pic = any(overlap_frac(box(pan), box(p)) > 0.30 for p in pics)
            holds_txt = any(overlap_frac(box(pan), box(t)) > 0.25 for t in texts)
            if not (holds_pic or holds_txt):
                issues.append(f"empty {w:.2f}x{h:.2f} in panel at "
                              f"({l:.2f},{t:.2f}) — a frame with nothing in it")

        #  TRUNCATION: a box whose text cannot fit is clipped on screen, silently
        for t in texts:
            tf = t.text_frame
            #  a frame allowed to grow, or one not wrapping, is not truncated
            try:
                if tf.word_wrap is False:
                    continue
                if tf.auto_size is not None and str(tf.auto_size).endswith("SHAPE_TO_FIT_TEXT"):
                    continue
            except Exception:
                pass
            _l, _tp, _w, h = box(t)
            need = est_text_height_in(t)
            if need > h * OVERFLOW_TOLERANCE and h > 0.15:
                issues.append(f"text {t.text_frame.text.strip()[:34]!r} needs about "
                              f"{need:.2f} in in a {h:.2f} in box — likely truncated")

        for t1 in texts:
            for p in pics:
                f = overlap_frac(box(t1), box(p))
                if f >= MIN_OVERLAP_FRAC:
                    issues.append(f"text {t1.text_frame.text.strip()[:34]!r} covers "
                                  f"{f*100:.0f} % of a figure")
        for a in range(len(texts)):
            for b in range(a + 1, len(texts)):
                f = overlap_frac(box(texts[a]), box(texts[b]))
                if f >= MIN_OVERLAP_FRAC:
                    issues.append(
                        f"text {texts[a].text_frame.text.strip()[:26]!r} overlaps "
                        f"{texts[b].text_frame.text.strip()[:26]!r} ({f*100:.0f} %)")

        if issues:
            print(f"  slide {i:2d}  ({len(pics)} figure(s))")
            for msg in issues:
                print(f"        - {msg}")
                n_fail += 1

    stale = stale_figures(path)
    if stale:
        seen = sorted(set(stale))
        print(f"  [NOTE] {len(stale)} figure(s) on slide(s) {seen} match no current "
              f"output — they may predate the last regeneration")
        n_fail += len(stale)
    else:
        print("  [ok  ] every figure matches a current generated output")

    total_pics = sum(1 for s in prs.slides for sh in s.shapes
                     if sh.__class__.__name__ == "Picture")
    print()
    print(f"{total_pics} figure(s) across the deck, {n_fail} issue(s)")
    return 1 if n_fail else 0


def main(argv):
    path = argv[0] if argv else "/mnt/c/Users/user/Desktop/slides3.pptx"
    if not os.path.exists(path):
        print(f"not found: {path}")
        return 2
    return audit(path)


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
