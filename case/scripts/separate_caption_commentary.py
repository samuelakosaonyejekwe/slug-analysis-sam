#!/usr/bin/env python3
"""Push a commentary block clear of the caption above it.

Moving figures between slides changes what sits under what: a caption that used to
sit beside a figure can end up overlapping a commentary block that never moved.

This resolves ONLY that pair — a short caption against a long commentary block — and
nothing else. An earlier, general "separate every overlapping text" pass made the
deck far worse, because a metric tile's value and its label overlap the tile card by
design and it dutifully pulled them apart. Overlap is not by itself a defect; this
kind is.

    python3 separate_caption_commentary.py [deck.pptx]
"""
import sys

from pptx import Presentation
from pptx.util import Emu

DECK = "/mnt/c/Users/user/Desktop/slides3.pptx"
CAPTION_MAX = 240         # longer than this is commentary, not a caption
FOOTER = 6.92
GAP = 0.08
MIN_H = 0.55
EMU = 914400


def main(argv):
    deck = argv[0] if argv else DECK
    prs = Presentation(deck)
    fixed = 0
    for i, slide in enumerate(prs.slides, 1):
        caps, comms = [], []
        for sh in slide.shapes:
            if sh.__class__.__name__ == "Picture" or not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if not t:
                continue
            try:
                g = (sh, Emu(sh.left).inches, Emu(sh.top).inches,
                     Emu(sh.width).inches, Emu(sh.height).inches)
            except Exception:
                continue
            if g[2] > FOOTER or g[2] + g[4] < 1.4:
                continue                       # footer, page number, header
            (caps if len(t) < CAPTION_MAX else comms).append(g)
        for c in caps:
            for m in comms:
                ox = min(c[1] + c[3], m[1] + m[3]) - max(c[1], m[1])
                oy = min(c[2] + c[4], m[2] + m[4]) - max(c[2], m[2])
                if ox <= 0.02 or oy <= 0.02:
                    continue
                want = c[2] + c[4] + GAP
                room = FOOTER - want
                if room < MIN_H:
                    continue
                m[0].top = int(round(want * EMU))
                if room < m[4]:
                    m[0].height = int(round(room * EMU))
                fixed += 1
                print(f"  slide {i:>2}  commentary moved below "
                      f"'{c[0].text_frame.text.strip()[:34]}'")
    prs.save(deck)
    print(f"\n  {fixed} caption/commentary overlap(s) resolved")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
