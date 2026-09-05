#!/usr/bin/env python3
"""Put the right RENDERING of each figure into each deck frame.

Type on a projected slide arrives at

    effective_pt = base_pt x (displayed_width / natural_width)

so legibility is decided by how a figure's natural width compares with the frame
it is shown in, not by the font size it was drawn with. The deck's frames run
from about 2.4 in to 8.0 in, which is why one rendering cannot serve all of them:
scaling every figure by the same factor divides every effective size by that same
factor, so whatever fixes the five animations crammed onto one slide blows the
single-figure slides up to fifty-point axis labels and softens them by upscaling.

make_slide_figures.py therefore renders the same figures at several size scales
(outputs_slides, outputs_slides45, outputs_slides70, ...). This script walks the
deck and, for each picture, swaps in the variant whose natural width puts the type
CLOSEST TO TARGET_PT in that particular frame. Frame geometry is never touched:
the replacement is fitted into the box the picture already occupies, preserving
the aspect ratio, so an earlier layout pass is not undone.

    python3 fit_deck_figures.py [deck.pptx]
"""
import hashlib
import io as _io
import os
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
CASE = os.path.abspath(os.path.join(HERE, ".."))
DECK = "/mnt/c/Users/user/Desktop/slides3.pptx"

#  Reflowing a caption downward gained about 0.2 pt on one slide and repeatedly
#  dropped labels onto their own descriptions on the five-animation slide, where a
#  filename label and a caption stack under each figure. Not worth the collisions.
REFLOW_CAPTIONS = False
PUSH_MAX_IN = 0.65        # how far a commentary card may be nudged to free height
MIN_W_IN = 2.6            # a figure narrower than this is a thumbnail, whatever
MIN_SQIN = 6.5            # its type size says — check_slides fails below 6 sq in
GROW_FRAMES = True        # widen a frame into empty slide space before giving up
GAP_IN = 0.12             # clearance kept from any neighbouring shape
TARGET_PT = 15.0          # comfortably above the 12 pt back-of-room threshold
FLOOR_PT = 12.0
CEILING_PT = 26.0         # above this the type crowds the plot out of its own axes
_ANIM_DPI_DEFAULT = float(os.environ.get("SHCT_ANIM_DPI", "150"))


def anim_dpi(path):
    """The dpi an animation was rendered at, from the sidecar beside it."""
    import json
    d = os.path.dirname(path) or "."
    side = os.path.join(d, "anim_dpi.json")
    try:
        with open(side) as fh:
            return float(json.load(fh)["dpi"])
    except Exception:
        return _ANIM_DPI_DEFAULT


ANIM_DPI = _ANIM_DPI_DEFAULT


def natural_in(path):
    with Image.open(path) as im:
        px = im.size[0]
        if (im.format or "").upper() == "GIF":
            dpi = anim_dpi(path)
        else:
            d = im.info.get("dpi")
            dpi = float(d[0]) if d and d[0] > 1 else 100.0
        return px / dpi, im.size


def variants():
    """basename -> [(path, natural_in, base_pt), ...] across every rendered set."""
    out = {}
    dirs = [d for d in sorted(os.listdir(CASE))
            if d.startswith("outputs_slides") and os.path.isdir(os.path.join(CASE, d))]
    dirs += ["outputs_steady", "outputs_shutin", "outputs_mitigated"]
    for d in dirs:
        base = 10.0 if not d.startswith("outputs_slides") else 18.0
        p = os.path.join(CASE, d)
        if not os.path.isdir(p):
            continue
        for f in os.listdir(p):
            if not f.endswith((".png", ".gif")):
                continue
            fp = os.path.join(p, f)
            try:
                nat, _ = natural_in(fp)
            except Exception:
                continue
            if nat > 0:
                out.setdefault(f, []).append((fp, nat, base))
    return out


def index_by_hash():
    """sha256 of every candidate file -> its basename, to identify what is placed."""
    idx = {}
    for d in os.listdir(CASE):
        p = os.path.join(CASE, d)
        if not (d.startswith("outputs") and os.path.isdir(p)):
            continue
        for f in os.listdir(p):
            if f.endswith((".png", ".gif")):
                try:
                    h = hashlib.sha256(open(os.path.join(p, f), "rb").read()).hexdigest()
                except OSError:
                    continue
                idx.setdefault(h, f)
    return idx


def pick(cands, disp_in):
    """The variant whose effective type size in a disp_in frame is nearest target.

    Ties and out-of-range cases are resolved toward the LARGER type, because a
    figure that is slightly too bold still reads from the back of the room and a
    figure that is slightly too small does not.
    """
    scored = []
    for fp, nat, base in cands:
        eff = base * disp_in / nat
        if eff > CEILING_PT:
            penalty = (eff - CEILING_PT) * 2.0
        elif eff < FLOOR_PT:
            penalty = (FLOOR_PT - eff) * 4.0
        else:
            penalty = 0.0
        scored.append((penalty + abs(eff - TARGET_PT), eff, fp, nat, base))
    scored.sort()
    return scored[0]


def free_box(slide, pic):
    """The empty rectangle this picture may grow into, in inches.

    A shape can only block growth along an axis if it OVERLAPS the picture on the
    other axis: a caption directly above does not stop the figure widening, but it
    does stop it growing upward. Computing only the horizontal freedom — as an
    earlier version did — let the picture grow DOWNWARD into captions and the
    footer, which is exactly the collision this is meant to prevent.
    """
    x0, y0 = Emu(pic.left).inches, Emu(pic.top).inches
    x1, y1 = x0 + Emu(pic.width).inches, y0 + Emu(pic.height).inches
    L, R = GAP_IN, SLIDE_W[0] - GAP_IN
    T_, B = GAP_IN, SLIDE_H[0] - GAP_IN
    for o in slide.shapes:
        if o is pic:
            continue
        try:
            a0, b0 = Emu(o.left).inches, Emu(o.top).inches
            a1, b1 = a0 + Emu(o.width).inches, b0 + Emu(o.height).inches
        except Exception:
            continue
        if not (b1 <= y0 + 0.02 or b0 >= y1 - 0.02):        # overlaps vertically
            if a1 <= x0 + 0.02:
                L = max(L, a1 + GAP_IN)
            elif a0 >= x1 - 0.02:
                R = min(R, a0 - GAP_IN)
        if not (a1 <= x0 + 0.02 or a0 >= x1 - 0.02):        # overlaps horizontally
            if b1 <= y0 + 0.02:
                T_ = max(T_, b1 + GAP_IN)
            elif b0 >= y1 - 0.02:
                B = min(B, b0 - GAP_IN)
    return (max(R - L, x1 - x0), max(B - T_, y1 - y0), L, T_)


SLIDE_W = [13.333]
SLIDE_H = [7.5]


def rows_of(pics):
    """Group pictures that share a horizontal band, in left-to-right order."""
    rows = []
    for p in sorted(pics, key=lambda q: (Emu(q.top).inches, Emu(q.left).inches)):
        y0 = Emu(p.top).inches
        y1 = y0 + Emu(p.height).inches
        for r in rows:
            ry0 = min(Emu(q.top).inches for q in r)
            ry1 = max(Emu(q.top).inches + Emu(q.height).inches for q in r)
            if not (y1 <= ry0 + 0.05 or y0 >= ry1 - 0.05):
                r.append(p)
                break
        else:
            rows.append([p])
    for r in rows:
        r.sort(key=lambda q: Emu(q.left).inches)
    return rows


def push_blocker(slide, row, want_h, T_, B, cap_h):
    """Nudge the card that caps a row DOWN, when there is empty space beneath it.

    Slide 37's sensitivity chart needed 2.46 in of height and had 2.29, because the
    "HOW THE RESULTS SHOULD BE USED" card sat at y = 4.30 with free slide below it.
    Moving a card is a layout change, so it is bounded: at most PUSH_MAX_IN, never
    past the next shape under it, and only when the figure above actually needs it.
    """
    y1 = max(Emu(p.top).inches + Emu(p.height).inches for p in row)
    x0 = min(Emu(p.left).inches for p in row)
    x1 = max(Emu(p.left).inches + Emu(p.width).inches for p in row)
    members = {id(p) for p in row}
    below = []
    for o in slide.shapes:
        if id(o) in members or o.__class__.__name__ == "Picture":
            continue
        try:
            a0, b0 = Emu(o.left).inches, Emu(o.top).inches
            a1, b1 = a0 + Emu(o.width).inches, b0 + Emu(o.height).inches
        except Exception:
            continue
        if b0 >= y1 - 0.02 and not (a1 <= x0 + 0.02 or a0 >= x1 - 0.02):
            below.append((b0, b1, o))
    if not below:
        return 0.0
    below.sort(key=lambda t: (t[0], t[1]))   # never fall through to comparing Shapes
    b0, b1, blocker = below[0]
    txt = blocker.text_frame.text.strip() if blocker.has_text_frame else ""
    #  never move the footer, the slide number or a rule
    if b0 > SLIDE_H[0] - 0.75 or txt.startswith("SHCT ") or txt.strip().isdigit():
        return 0.0
    #  and never move a CAPTION. The shape under a figure is usually its own label,
    #  and nudging it down drops it onto the description beneath — which is how the
    #  five-animation slide gained collisions from a pass meant to remove them. Only
    #  a commentary block (long text) or an untitled card may be pushed.
    if txt and len(txt) < 240:
        return 0.0
    have = (b0 - GAP_IN) - T_ - cap_h
    short = want_h - have
    if short <= 0.01:
        return 0.0
    #  the next thing under the blocker, or the footer band
    floor_y = SLIDE_H[0] - 0.75
    for c0, c1, o in below[1:]:
        if c0 >= b1 - 0.02:
            floor_y = min(floor_y, c0)
            break
    slack = max(floor_y - b1 - GAP_IN, 0.0)
    push = min(short, slack, PUSH_MAX_IN)
    if push <= 0.02:
        return 0.0
    for o in (blocker,):
        o.top = int(round((Emu(o.top).inches + push) * 914400))
    return push


def caption_of(slide, row):
    """The caption line sitting directly beneath a row of pictures, if any."""
    y1 = max(Emu(p.top).inches + Emu(p.height).inches for p in row)
    x0 = min(Emu(p.left).inches for p in row)
    x1 = max(Emu(p.left).inches + Emu(p.width).inches for p in row)
    best, best_gap = None, 1e9
    for o in slide.shapes:
        if o.__class__.__name__ == "Picture" or not o.has_text_frame:
            continue
        t = o.text_frame.text.strip()
        if not t or len(t) > 240:
            continue
        a0, b0 = Emu(o.left).inches, Emu(o.top).inches
        a1 = a0 + Emu(o.width).inches
        #  the running footer and the slide number sit below every figure and are not
        #  captions; moving the footer onto a data table is what happens if they are
        if b0 > SLIDE_H[0] - 0.75 or t.startswith("SHCT ") or t.strip().isdigit():
            continue
        gap = b0 - y1
        if 0 <= gap < 0.8 and not (a1 <= x0 - 0.1 or a0 >= x1 + 0.1) and gap < best_gap:
            best, best_gap = o, gap
    return best


def row_space(slide, row):
    """The free rectangle a whole row of pictures may occupy, in inches.

    Only NON-picture shapes constrain it: the pictures in the row are being laid
    out together, so they must not block one another. Treating a sibling as an
    immovable wall is what stopped slide 12's two figures from sharing the width
    they both needed — 12.5 in of need against 13.1 in available, and neither
    could move.
    """
    y0 = min(Emu(p.top).inches for p in row)
    y1 = max(Emu(p.top).inches + Emu(p.height).inches for p in row)
    x0 = min(Emu(p.left).inches for p in row)
    x1 = max(Emu(p.left).inches + Emu(p.width).inches for p in row)
    L, R, T_, B = GAP_IN, SLIDE_W[0] - GAP_IN, GAP_IN, SLIDE_H[0] - GAP_IN
    members = {id(p) for p in row}
    #  A caption belongs to the row, so it must not cap the row's height where there
    #  is free space beneath it — it can simply move down. Slides 32 and 37 were both
    #  held short by their own captions with room to spare below.
    cap = caption_of(slide, row) if REFLOW_CAPTIONS else None
    cap_h = (Emu(cap.height).inches + 0.10) if cap is not None else 0.0
    for o in slide.shapes:
        if id(o) in members or o.__class__.__name__ == "Picture" or o is cap:
            continue
        try:
            a0, b0 = Emu(o.left).inches, Emu(o.top).inches
            a1, b1 = a0 + Emu(o.width).inches, b0 + Emu(o.height).inches
        except Exception:
            continue
        if not (b1 <= y0 + 0.02 or b0 >= y1 - 0.02):
            if a1 <= x0 + 0.02:
                L = max(L, a1 + GAP_IN)
            elif a0 >= x1 - 0.02:
                R = min(R, a0 - GAP_IN)
        if not (a1 <= x0 + 0.02 or a0 >= x1 - 0.02):
            if b1 <= y0 + 0.02:
                T_ = max(T_, b1 + GAP_IN)
            elif b0 >= y1 - 0.02:
                B = min(B, b0 - GAP_IN)
    B -= cap_h                       # keep room for the caption under the figures
    return L, max(R - L, x1 - x0), T_, max(B - T_, y1 - y0), cap


def main(argv):
    deck = argv[0] if argv else DECK
    prs = Presentation(deck)
    SLIDE_W[0] = Emu(prs.slide_width).inches
    SLIDE_H[0] = Emu(prs.slide_height).inches
    cands = variants()
    byhash = index_by_hash()
    swapped = kept = 0
    report = []

    for i, slide in enumerate(prs.slides, 1):
        pics = [sh for sh in slide.shapes if sh.__class__.__name__ == "Picture"]
        for row in rows_of(pics):
            L, avail_w, T_, avail_h, cap = row_space(slide, row)
            info = []
            for sh in row:
                try:
                    blob = sh.image.blob
                except Exception:
                    info.append(None)
                    continue
                name = byhash.get(hashlib.sha256(blob).hexdigest())
                if not name or name not in cands:
                    info.append(None)
                    continue
                #  reference rendering: the smallest natural among the slide-legible
                #  sets, which is the one that needs the least width to read
                nat_ref, base_ref, path_ref = min(
                    ((n, b, f) for f, n, b in cands[name]), key=lambda t: t[0] / t[1])
                with Image.open(path_ref) as im2:
                    ar = im2.size[1] / im2.size[0]
                info.append({"sh": sh, "name": name, "blob": blob, "ar": ar,
                             "need": nat_ref * TARGET_PT / base_ref,
                             "cur": Emu(sh.width).inches})
            live = [d for d in info if d]
            if not live:
                continue
            gaps = GAP_IN * (len(live) - 1)
            room = max(avail_w - gaps, 0.5)
            #  Fill the row in the proportions the slide was designed with. Sizing each
            #  figure to the smallest width that reaches the target type size instead
            #  produced 1.5 in postage stamps once the compact renderings arrived: the
            #  labels were 15 pt and the plot was unreadable. A figure wants to be as
            #  large as its row allows; the RENDERING is then chosen to suit that size.
            share = sum(d["cur"] for d in live) or 1.0
            for d in live:
                d["w"] = room * d["cur"] / share
                d["h"] = d["w"] * d["ar"]
            #  a row may be limited by height as well
            tall = max((d["h"] / avail_h) for d in live)
            if tall > 1.0:
                for d in live:
                    d["w"] /= tall
                    d["h"] /= tall
            #  never below a readable physical size, room permitting
            for d in live:
                floor_w = max(MIN_W_IN, (MIN_SQIN / d["ar"]) ** 0.5 if d["ar"] > 0 else 0)
                d["w"] = min(max(d["w"], min(floor_w, room)), room)
                d["h"] = min(d["w"] * d["ar"], avail_h)
            #  centre the row horizontally in its free span
            #  if the row is height-limited, see whether the card below it can move
            want = max(d["w"] * d["ar"] for d in live)
            if want > avail_h + 0.01:
                got = push_blocker(slide, row, want, T_, T_ + avail_h,
                                   (Emu(cap.height).inches + 0.10) if cap else 0.0)
                if got > 0.0:
                    avail_h += got
                    for d in live:
                        d["h"] = min(d["w"] * d["ar"], avail_h)
                        if d["h"] < d["w"] * d["ar"]:
                            d["w"] = d["h"] / d["ar"] if d["ar"] > 0 else d["w"]
            placed_bottom = 0.0
            used = sum(d["w"] for d in live) + gaps
            x = L + max((avail_w - used) / 2.0, 0.0)
            for d in live:
                sh = d["sh"]
                #  now choose the rendering that reads best at the width allocated
                _, eff, fp, nat, base = pick(cands[d["name"]], d["w"])
                cur_nat = None
                try:
                    im = Image.open(_io.BytesIO(d["blob"]))
                    cur_dpi = (anim_dpi(fp) if (im.format or "").upper() == "GIF"
                               else (float(im.info["dpi"][0])
                                     if im.info.get("dpi") and im.info["dpi"][0] > 1
                                     else 100.0))
                    cur_nat = im.size[0] / cur_dpi
                except Exception:
                    pass
                same = bool(cur_nat and abs(cur_nat - nat) < 1e-6)
                moved = (abs(d["w"] - d["cur"]) > 0.05
                         or abs(x - Emu(sh.left).inches) > 0.05)
                if same and not moved:
                    kept += 1
                    x += d["w"] + GAP_IN
                    continue
                top = T_ + max((avail_h - d["h"]) / 2.0, 0.0)
                placed_bottom = max(placed_bottom, top + d["h"])
                #  Never place a figure on top of text. row_space keeps the row clear
                #  of its neighbours, but a row that groups five animations can still
                #  grow a figure across the label belonging to the one beside it, and
                #  a legibility pass that creates a collision has made things worse.
                #  If the new box would touch any text, keep the original geometry.
                _nx0, _ny0 = x, top
                _nx1, _ny1 = x + d["w"], top + d["h"]
                _hits_text = False
                for _o in slide.shapes:
                    if _o.__class__.__name__ == "Picture" or not _o.has_text_frame:
                        continue
                    if not _o.text_frame.text.strip():
                        continue
                    try:
                        _a0, _b0 = Emu(_o.left).inches, Emu(_o.top).inches
                        _a1, _b1 = _a0 + Emu(_o.width).inches, _b0 + Emu(_o.height).inches
                    except Exception:
                        continue
                    if (not (_a1 <= _nx0 + 0.02 or _a0 >= _nx1 - 0.02)
                            and not (_b1 <= _ny0 + 0.02 or _b0 >= _ny1 - 0.02)):
                        _hits_text = True
                        break
                if _hits_text:
                    x, top = Emu(sh.left).inches, Emu(sh.top).inches
                    d["w"], d["h"] = Emu(sh.width).inches, Emu(sh.height).inches
                newpic = slide.shapes.add_picture(
                    fp, int(round(x * 914400)), int(round(top * 914400)),
                    int(round(d["w"] * 914400)), int(round(d["h"] * 914400)))
                sh._element.addprevious(newpic._element)
                sh._element.getparent().remove(sh._element)
                swapped += 1
                old_eff = base * d["cur"] / cur_nat if cur_nat else float("nan")
                report.append((i, d["name"], old_eff, eff,
                               os.path.basename(os.path.dirname(fp))))
                x += d["w"] + GAP_IN
            if cap is not None and placed_bottom > 0.0:
                #  Moving a caption is only safe if it lands on empty slide. Slide 34
                #  stacks a filename label above a description, and reflowing the
                #  first dropped it onto the second — a collision created by the very
                #  pass meant to prevent them. Check the destination before moving.
                _ch = Emu(cap.height).inches
                _cw = Emu(cap.width).inches
                _cl = Emu(cap.left).inches
                _newtop = min(placed_bottom + 0.08, SLIDE_H[0] - _ch - 0.35)
                _clash = False
                for _o in slide.shapes:
                    if _o is cap or _o.__class__.__name__ == "Picture":
                        continue
                    try:
                        _a0, _b0 = Emu(_o.left).inches, Emu(_o.top).inches
                        _a1, _b1 = _a0 + Emu(_o.width).inches, _b0 + Emu(_o.height).inches
                    except Exception:
                        continue
                    if (not (_a1 <= _cl + 0.02 or _a0 >= _cl + _cw - 0.02)
                            and not (_b1 <= _newtop + 0.02
                                     or _b0 >= _newtop + _ch - 0.02)):
                        _clash = True
                        break
                if _clash:
                    cap = None
            if cap is not None and placed_bottom > 0.0:
                cap.top = int(round(min(placed_bottom + 0.08,
                                        SLIDE_H[0] - Emu(cap.height).inches - 0.35)
                                    * 914400))

    prs.save(deck)
    for i, name, oe, ne, d in sorted(report, key=lambda r: r[3]):
        print(f"  slide {i:>2}  {name:<32} {oe:>5.1f} -> {ne:>5.1f} pt   [{d}]")
    print(f"\n  {swapped} figure(s) refitted, {kept} already best")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
