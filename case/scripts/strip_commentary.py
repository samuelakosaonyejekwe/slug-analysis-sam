#!/usr/bin/env python3
"""Move a slide's commentary into its speaker notes when the figure needs the room.

The same trade that rescued the two space-time slides, applied wherever it is still
needed: a figure set beside a block of commentary gets about a third of the slide,
and for a dense multi-panel plot that puts its labels below the back-of-room
threshold. Moving the commentary into the SPEAKER NOTES gives the figure the whole
width without losing a word — the presenter still has the text, the audience gets a
figure they can read.

Only slides that still fall short are touched, and only the commentary is moved.
Captions stay with their figures, metric tiles stay on the slide (they are the
numbers being pointed at), and the header and footer are untouched. Run it after
fit_deck_figures.py has placed everything, then fit again so the figures expand
into the space released.

    python3 strip_commentary.py [deck.pptx] [--threshold PT]
"""
import hashlib
import io as _io
import os
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
CASE = os.path.abspath(os.path.join(HERE, ".."))
DECK = "/mnt/c/Users/user/Desktop/slides3.pptx"

THRESHOLD_PT = 12.0
MIN_SQIN = 6.0
TOP_BAND, BOT_BAND = 1.56, 6.88
TILE_BAND = 1.45          # a block ending this close to the top band is a metric tile
COMMENTARY_CHARS = 150    # a caption is short; this is the length side of the test
CAPTION_UNDER_IN = 0.9    # ...and a caption sits directly UNDER its figure
ANIM_DPI_DEFAULT = 200.0


def anim_dpi(path):
    import json
    try:
        with open(os.path.join(os.path.dirname(path), "anim_dpi.json")) as fh:
            return float(json.load(fh)["dpi"])
    except Exception:
        return ANIM_DPI_DEFAULT


def index():
    idx = {}
    for d in os.listdir(CASE):
        p = os.path.join(CASE, d)
        if not (d.startswith("outputs") and os.path.isdir(p)):
            continue
        base = 18.0 if d.startswith("outputs_slides") else 10.0
        for f in os.listdir(p):
            if not f.endswith((".png", ".gif")):
                continue
            fp = os.path.join(p, f)
            try:
                h = hashlib.sha256(open(fp, "rb").read()).hexdigest()
            except OSError:
                continue
            idx.setdefault(h, (f, base, fp))
    return idx


def effective_pt(sh, base, src):
    try:
        im = Image.open(_io.BytesIO(sh.image.blob))
        if (im.format or "").upper() == "GIF":
            dpi = anim_dpi(src)
        else:
            d = im.info.get("dpi")
            dpi = float(d[0]) if d and d[0] > 1 else 100.0
        nat = im.size[0] / dpi
        return base * Emu(sh.width).inches / nat if nat else 0.0
    except Exception:
        return 0.0


def main(argv):
    deck, thr = DECK, THRESHOLD_PT
    args = list(argv)
    if "--threshold" in args:
        k = args.index("--threshold")
        thr = float(args[k + 1])
        del args[k:k + 2]
    if args:
        deck = args[0]

    prs = Presentation(deck)
    idx = index()
    touched = 0

    for i, slide in enumerate(prs.slides, 1):
        pics = [s for s in slide.shapes if s.__class__.__name__ == "Picture"]
        if not pics:
            continue
        short = False
        for sh in pics:
            name, base, src = idx.get(
                hashlib.sha256(sh.image.blob).hexdigest(), (None, 18.0, ""))
            if name is None:
                continue
            area = Emu(sh.width).inches * Emu(sh.height).inches
            if effective_pt(sh, base, src) < thr or area < MIN_SQIN:
                short = True
                break
        if not short:
            continue

        moved, drop = [], []
        for sh in slide.shapes:
            if sh.__class__.__name__ == "Picture" or not sh.has_text_frame:
                continue
            txt = sh.text_frame.text.strip()
            if not txt:
                continue
            try:
                t, h = Emu(sh.top).inches, Emu(sh.height).inches
            except Exception:
                continue
            #  Judge chrome by where a block STARTS. Testing its bottom edge made a
            #  commentary block that runs to the footer rule look like the footer, and
            #  slide 14 was skipped entirely with its figure stuck at 4.82 in.
            if t < TOP_BAND - 0.05 or t > BOT_BAND:
                continue                                   # header / footer
            if t + h <= TOP_BAND + TILE_BAND:
                continue                                   # metric tile
            #  Length alone mis-sorts: slide 14 carries a 221-character block BESIDE
            #  the figure (commentary, and it was capping the figure at 5.75 in) while
            #  slide 45's caption runs to 189. What separates them is position — a
            #  caption sits directly under its figure and roughly shares its width.
            under = any(0 <= t - (Emu(p2.top).inches + Emu(p2.height).inches)
                        < CAPTION_UNDER_IN
                        and abs(Emu(p2.left).inches - Emu(sh.left).inches) < 1.9
                        for p2 in pics)
            if under or len(txt) < COMMENTARY_CHARS:
                continue                                   # a caption
            moved.append(txt)
            drop.append(sh)
        if not moved:
            continue

        tf = slide.notes_slide.notes_text_frame
        was = tf.text.strip()
        tf.text = (was + "\n\n" if was else "") + "\n\n".join(moved)
        for sh in drop:
            sh._element.getparent().remove(sh._element)
        #  the card behind the commentary, now empty, goes with it
        for sh in list(slide.shapes):
            if sh.__class__.__name__ == "Picture":
                continue
            if sh.has_text_frame and sh.text_frame.text.strip():
                continue
            try:
                t, h = Emu(sh.top).inches, Emu(sh.height).inches
                w = Emu(sh.width).inches
            except Exception:
                continue
            if TOP_BAND - 0.05 <= t and t + h <= BOT_BAND + 0.05 and w > 1.5 and h > 0.8:
                sh._element.getparent().remove(sh._element)
        touched += 1
        print(f"  slide {i:>2}  {len(moved)} commentary block(s), "
              f"{sum(len(m) for m in moved)} chars -> speaker notes")

    prs.save(deck)
    print(f"\n  {touched} slide(s) cleared for their figures")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
