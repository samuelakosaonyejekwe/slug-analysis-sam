#!/usr/bin/env python3
"""Stand the space-time map across its slide, and speak the commentary instead.

The six-panel space-time figure renders 12.97 in wide. Set beside a commentary
card it gets about a third of the slide and its labels land at 7-8 pt; across the
full width it reads at about 17 pt. Both slides carrying it have the width and
neither has the height to stack figure over commentary, so the commentary has to
leave the slide.

It is not deleted. The text moves into the slide's SPEAKER NOTES, where a presenter
still has every word and the audience gets a figure they can read. Metric tiles
stay on the slide — they are the numbers being pointed at.

    python3 restack_spacetime_slides.py [deck.pptx]
"""
import hashlib
import os
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
CASE = os.path.abspath(os.path.join(HERE, ".."))
DECK = "/mnt/c/Users/user/Desktop/slides3.pptx"

TARGET = "19_spacetime_fields.png"
MARGIN, TOP_BAND, BOT_BAND = 0.62, 1.56, 6.88
TILE_BAND = 1.45          # anything ending within this of the top band is a metric tile
GAP = 0.16
EMU = 914400


def widest(name):
    best = None
    for d in sorted(os.listdir(CASE)):
        p = os.path.join(CASE, d, name)
        if not (d.startswith("outputs_slides") and os.path.isfile(p)):
            continue
        try:
            with Image.open(p) as im:
                px, py = im.size
                dd = im.info.get("dpi")
                dpi = float(dd[0]) if dd and dd[0] > 1 else 100.0
        except Exception:
            continue
        ar = py / px
        if best is None or ar < best[2]:
            best = (p, px / dpi, ar)
    return best


def main(argv):
    deck = argv[0] if argv else DECK
    prs = Presentation(deck)
    full_w = Emu(prs.slide_width).inches - 2 * MARGIN
    wv = widest(TARGET)
    if wv is None:
        print(f"  no rendering of {TARGET} found")
        return 2
    path, nat, ar = wv

    idx = {}
    for d in os.listdir(CASE):
        p = os.path.join(CASE, d)
        if not (d.startswith("outputs") and os.path.isdir(p)):
            continue
        for f in os.listdir(p):
            if f.endswith((".png", ".gif")):
                try:
                    idx.setdefault(hashlib.sha256(
                        open(os.path.join(p, f), "rb").read()).hexdigest(), f)
                except OSError:
                    pass

    touched = 0
    for i, slide in enumerate(prs.slides, 1):
        pics = [s for s in slide.shapes if s.__class__.__name__ == "Picture"]
        if len(pics) != 1:
            continue
        pic = pics[0]
        if idx.get(hashlib.sha256(pic.image.blob).hexdigest()) != TARGET:
            continue

        tiles_bottom, moved_text, drop = TOP_BAND, [], []
        caption = None
        for sh in slide.shapes:
            if sh is pic:
                continue
            try:
                t, h = Emu(sh.top).inches, Emu(sh.height).inches
                l, w = Emu(sh.left).inches, Emu(sh.width).inches
            except Exception:
                continue
            if t < TOP_BAND - 0.05 or t + h > BOT_BAND + 0.05:
                continue                                  # header / footer chrome
            txt = sh.text_frame.text.strip() if sh.has_text_frame else ""
            if t + h <= TOP_BAND + TILE_BAND:
                tiles_bottom = max(tiles_bottom, t + h)    # a metric tile: keep it
                continue
            if txt and len(txt) < 240 and caption is None and w > 2.0:
                caption = sh                              # the figure's own caption
                continue
            if txt:
                moved_text.append(txt)
            drop.append(sh)

        #  the commentary is preserved where a presenter can still use it
        if moved_text:
            tf = slide.notes_slide.notes_text_frame
            existing = tf.text.strip()
            tf.text = (existing + "\n\n" if existing else "") + "\n\n".join(moved_text)
        for sh in drop:
            sh._element.getparent().remove(sh._element)

        top0 = tiles_bottom + (GAP if tiles_bottom > TOP_BAND else 0.0)
        fig_h = min(full_w * ar, BOT_BAND - top0 - 0.45)
        fig_w = fig_h / ar if ar > 0 else full_w
        fig_w = min(fig_w, full_w)
        left = MARGIN + (full_w - fig_w) / 2.0
        pic._element.getparent().remove(pic._element)
        slide.shapes.add_picture(path, int(left * EMU), int(top0 * EMU),
                                 int(fig_w * EMU), int(fig_h * EMU))
        if caption is not None:
            caption.left = int(MARGIN * EMU)
            caption.top = int((top0 + fig_h + 0.08) * EMU)
            caption.width = int(full_w * EMU)
        touched += 1
        print(f"  slide {i:>2}  {TARGET} across {fig_w:.2f} in -> "
              f"~{18.0 * fig_w / nat:.1f} pt; {len(moved_text)} commentary block(s) "
              f"moved to the speaker notes")

    prs.save(deck)
    print(f"\n  {touched} slide(s) restacked")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
