#!/usr/bin/env python3
# =============================================================================
#  check_docs.py — find superseded numbers in the manuscript, thesis and slides.
# -----------------------------------------------------------------------------
#  A generated figure can be checked for freshness by its modification time
#  (check_journal_artwork.py does that). A DOCUMENT cannot: prose carries the
#  results as text, and a paragraph written against an old run looks exactly like
#  one written against the current run. The only mechanical handle is the values
#  themselves, so this keeps a register of retired values and reports any that
#  are still in the text.
#
#  Two failure modes made hand-scanning unreliable, and both are handled here
#  because both actually happened during the v3.3.0 update:
#
#    NOTATION.  The same quantity appears as 1.51×10⁻³, 1.51x10-3, 1.5e-3 and
#               "0.15 %". A search for one of those is blind to the other three,
#               and a mass-conservation error that had been fixed survived in the
#               thesis conclusions purely because it was written in ASCII while
#               every scan used unicode. Each entry below therefore carries all
#               the spellings the value has actually been seen in, and the text
#               is normalised (unicode superscripts and × folded to ASCII, thin
#               and non-breaking spaces to plain) before matching.
#
#    STRUCTURE. In a .docx table every cell is its own <w:p>. A value sitting
#               alone in a cell — "6.605" — is invisible to any search that joins
#               a paragraph's runs, because the label is in a different cell
#               entirely. In a .pptx the same is true of KPI tiles. Both are
#               therefore walked element by element, and every cell is treated as
#               its own searchable unit rather than as part of a sentence.
#
#      python3 check_docs.py [file ...]
#
#  With no arguments it checks the documents of this project. Exit status is 1 if
#  any retired value is found.
# =============================================================================
import os
import re
import sys
import unicodedata

W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"

#  (label, [spellings], what it should read now)
#  Every entry here was a real defect found in this project's documents.
RETIRED = [
    ("liquid mass-conservation error (was a bug, now round-off)",
     [r"1\.51\s*[x×*]\s*10\s*[-−⁻]\s*3", r"1\.5\s*[x×*]\s*10\s*[-−⁻]\s*3",
      r"1\.51e-0?3", r"0\.15\s*%"], "1.63e-13"),
    ("gas mass-conservation error",
     [r"1\.77\s*[x×*]\s*10\s*[-−⁻]\s*15", r"1\.8\s*[x×*]\s*10\s*[-−⁻]\s*15"], "1.68e-16"),
    ("P50 time-to-plug, as-operated", [r"\b2\.78\s*h", r"\b2\.77\s*h", r"\b2\.8\s*h"], "3.18 h"),
    ("P10/P50/P90 band, as-operated", [r"2\.13\s*/\s*2\.78", r"2\.13\s*h"], "2.45/3.18/4.27 h"),
    ("total pressure drop, as-operated", [r"\b113\.8\b"], "135.7 bar"),
    ("max subcooling, as-operated", [r"\b20\.9\s*°?C", r"\b21\s*°C\b"], "17.6 C"),
    ("P90 design subcooling, as-operated", [r"\b23\.7\s*°?C"], "20.4 C"),
    ("intermittent/slug fraction", [r"\b0\.814\b"], "0.789"),
    ("mean slug length", [r"\b15\.6\s*m\b"], "24.3 m"),
    ("max slug length", [r"\b36\.6\s*m\b"], "37.1 m"),
    ("hydrate mass formed, as-operated", [r"9\.76\s*[x×*]\s*10", r"\b9\.76\b"], "8.10e6 kg"),
    ("arrival temperature", [r"\b5\.3\s*°?C"], "13.9 C"),
    ("MEG dose, as-operated", [r"\b59\.7\s*wt"], "60 wt% (ceiling)"),
    ("MEG rate, as-operated", [r"94[,\s]470"], "80,095 L/h"),
    ("under-inhibited length, as-operated", [r"\b24\.2\s*km"], "25.1 km"),
    ("mitigated plug probability", [r"\b25\s*%\s*(?:residual|plug)", r"to\s*25\s*%",
                                    r"falls to 25", r"Plug probability 0\.25"], "zero"),
    ("mitigated max subcooling", [r"\b6\.605\b", r"\b6\.6\s*°?C"], "9.14 C"),
    ("mitigated P90 subcooling", [r"\b6\.641\b"], "10.3 C"),
    ("mitigated MEG dose", [r"\b30\.5\s*wt", r"\b30\.52\b"], "39.1 wt%"),
    ("mitigated MEG rate", [r"28[,\s]038"], "40,967 L/h"),
    ("mitigated under-inhibited length", [r"\b5\.49\s*km", r"\b5\.5\s*km\s*under"], "3.2 km"),
    ("mitigated peak Phi_SH", [r"\b3815\b"], "3800"),
    ("mitigated time-to-plug band", [r"3\.02\s*/\s*3\.93", r"\b3\.93\s*h"], "none — nothing plugs"),
    ("shut-in plug probability", [r"\b92\s*%"], "100%"),
    ("shut-in max subcooling", [r"\b28\.6\s*°?C"], "28.8 C"),
    ("sustained Phi_SH, as-operated", [r"=\s*4\.15\b", r"\b4\.15\s*(as operated|over)"], "2593"),
    ("super-critical length, as-operated", [r"\b16\.9\s*km", r"\b16\.5\s*km"], "22.9 km"),
    ("sensitivity: time-to-plug spread", [r"17\.5\s*h", r"factor of (?:roughly )?(?:twenty|20)\b"],
     "4.3 h to 1.0 h, factor of nine"),
    ("sensitivity: gate-saturated fraction", [r"gate[- ]saturat", r"\b44\s*%"], "24 % above Phi_crit"),
    ("the removed Phi_SH gate", [r"Φ_SH[- ]gated", r"ΦSH[- ]gated", r"gated by the coupling",
                                 r"f_wall opens", r"gating of the coupling",
                                 r"gates the wall-capture", r"clip\(Φ", r"min\(max\(Φ",
                                 r"consolidation only above criticality"], "rates compete; no gate"),
    #  --- superseded by the wall-area correction and the late-life scenario ---
    #  The wall growth law used the GAS-LIQUID interfacial area, which vanishes as a
    #  line fills with liquid; corrected to the wall area it agrees with Qin's (2020)
    #  measured film growth. Phi_SH, formed from the same rate, fell by ~13x, and the
    #  as-operated case moved to late-life conditions (70 % water cut, 0.6x rate) so
    #  that it exercises the threshold instead of sitting orders of magnitude above it.
    ("sustained Phi_SH, as-operated", [r"\b2593\b", r"=\s*2593"], "1.06"),
    ("peak Phi_SH, as-operated", [r"\b6123\b", r"6\.12\s*[x×]\s*10", r"\b6627\b", r"6\.12x10"], "1.95"),
    ("mitigated peak Phi_SH", [r"\b3800\b", r"3\.80\s*[x×]\s*10"], "1.00"),
    ("super-critical length, as-operated", [r"\b22\.9\s*km", r"\b20\.6\s*km"], "1.37 km"),
    ("P50 time-to-plug, as-operated", [r"\b3\.18\s*h", r"\b3\.2\s*h"], "3.72 h"),
    ("P10/P90 band, as-operated", [r"2\.45\s*/\s*3\.18", r"\b2\.45\s*h", r"\b4\.27\s*h"],
     "3.02/3.72/5.96 h"),
    ("max subcooling, as-operated", [r"\b17\.6\s*°?C"], "24.4 C"),
    ("MEG dose, as-operated", [r"\b55\.7\s*wt", r"\b56\s*wt%"], "60 wt% (ceiling)"),
    ("under-inhibited length, as-operated", [r"\b25\.1\s*km"], "26.5 km"),
    ("water cut", [r"\b35\s*%\s*water", r"water cut.{0,12}\b35\b"], "70 %"),
    ("mitigated plug probability (now non-zero)",
     [r"plug probability (?:to|falls to|drops? from 100% to) zero",
      r"no realization plugs"], "0.33"),
    ("mitigated peak deposit",
     [r"eliminates the deposit \(10\.2", r"peak deposit to 10\.2",
      r"10\.2 mm against full bore"], "16.4 mm"),
    ("hydrate mass, as-operated", [r"8\.10\s*[x×]\s*10", r"\b8\.1e\+?06"], "6.77e6 kg"),
    ("total dP, as-operated", [r"\b135\.7\b"], "80.4 bar"),
    ("arrival temperature", [r"\b13\.9\s*°?C"], "5.5 C"),
    ("mean slug length", [r"\b24\.3\s*m\b"], "26.9 m"),
    ("max slug length", [r"\b37\.1\s*m\b"], "59.4 m"),
    ("intermittent fraction", [r"\b0\.789\b"], "0.727"),
    ("the gas-liquid area in the wall growth law",
     [r"wall growth.{0,40}interfacial area", r"a_i.{0,20}wall growth",
      r"deposit.{0,30}driven by.{0,40}interfacial area"], "wall area 4/D*alpha_l*wf"),

    ("the threshold stated as assumed", [r"unity by construction(?!\s*—)",
                                         r"threshold of Φ_SH is unity",
                                         r"Φ_SH\s*[<>]\s*1(?![\d.])",
                                         r"ΦSH\s*=\s*1\s*(?:contour|criterion|threshold)"],
     "derived Phi_crit = 1.08"),
]

#  Deliberate references to a superseded value — a passage that says what a number
#  USED to be, and why it changed, is not a stale number. Recording the correction
#  is the honest thing to do, so the check must not punish it. Matched against the
#  normalised text of the same unit.
ALLOW = [
    "rather than set to unity",              # the abstract, stating what was NOT done
    "an earlier revision of the solver reported",   # Ch.6, recording the fixed defect
    "an earlier formulation gated",          # Ch.5, recording the removed gate
    "a formulation that gates wall capture", # paper 3.3, the counterfactual
    "the earlier gated formulation",         # paper 5, the discriminating comparison
    "which made the criticality threshold an input",
]

TRANS = str.maketrans({"⁻": "-", "−": "-", "×": "x", " ": " ", " ": " ",
                       "⁰": "0", "¹": "1", "²": "2", "³": "3", "⁴": "4",
                       "⁵": "5", "⁶": "6", "⁷": "7", "⁸": "8", "⁹": "9"})


def norm(s):
    """Fold the notation variants a number is written in onto one spelling."""
    return unicodedata.normalize("NFKC", s).translate(TRANS)


def units_docx(path):
    """Every independently-searchable unit of text: one per <w:p>.

    A table cell is its own <w:p>, so a bare value in a cell is yielded on its
    own rather than glued to a neighbouring sentence — which is the only way a
    lone "6.605" in a results table can be found at all.
    """
    from docx import Document
    d = Document(path)
    parts = [d.element]
    for sec in d.sections:
        for hf in (sec.header, sec.footer, sec.first_page_header, sec.first_page_footer):
            try:
                parts.append(hf._element)
            except Exception:
                pass
    n = 0
    for part in parts:
        for p in part.iter(f"{W}p"):
            #  .//w:t catches runs inside hyperlinks and fields, which
            #  python-docx's paragraph.runs silently omits
            txt = "".join(t.text or "" for t in p.findall(f".//{W}t"))
            n += 1
            if txt.strip():
                yield f"para {n}", txt


def units_pptx(path):
    """One unit per paragraph of every text frame, plus every table cell."""
    from pptx import Presentation
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if sh.has_text_frame:
                for j, para in enumerate(sh.text_frame.paragraphs):
                    txt = "".join(r.text for r in para.runs)
                    if txt.strip():
                        yield f"slide {i} para {j}", txt
            if getattr(sh, "has_table", False):
                for r, row in enumerate(sh.table.rows):
                    for c, cell in enumerate(row.cells):
                        if cell.text.strip():
                            yield f"slide {i} cell[{r}][{c}]", cell.text


def check_pdf_freshness(paths):
    """A .pdf sitting beside a .docx must not be older than it.

    This exists because the check above did not catch the thing that mattered most.
    It scanned the .docx and the .pptx and reported every document consistent, while
    paper5.pdf, paper5_typeset.pdf and slugs1-thesis.pdf were two days stale — the
    Word files had been corrected and the PDFs never regenerated. The PDF is the
    artefact that gets submitted and read, so checking only its source is checking
    the wrong end of the pipeline.

    Modification time is the handle, as it is for figures: a PDF older than the
    document it was exported from cannot contain that document's current text.
    """
    out = []
    seen = set()
    for src in paths:
        if not src.lower().endswith(".docx"):
            continue
        pdf = src[:-5] + ".pdf"
        if pdf in seen or not os.path.exists(pdf):
            continue
        seen.add(pdf)
        t_src, t_pdf = os.path.getmtime(src), os.path.getmtime(pdf)
        if t_pdf < t_src - 1.0:
            out.append((pdf, (t_src - t_pdf) / 3600.0))
    return out


def sibling_docs(paths):
    """Every .docx in the directories being checked, so a PDF is not missed merely
    because its source was not named on the command line."""
    found = list(paths)
    for d in {os.path.dirname(os.path.abspath(p)) for p in paths}:
        try:
            for f in os.listdir(d):
                if (f.endswith(".docx") and not f.startswith("~$")
                        and ".pre-" not in f):
                    full = os.path.join(d, f)
                    if full not in found:
                        found.append(full)
        except OSError:
            pass
    return found


def check(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx":
        units = units_docx
    elif ext == ".pptx":
        units = units_pptx
    else:
        print(f"  [skip] {os.path.basename(path)}: not a .docx or .pptx")
        return 0

    compiled = [(label, [re.compile(norm(p), re.I) for p in pats], now)
                for label, pats, now in RETIRED]
    hits = []
    for where, raw in units(path):
        txt = norm(raw)
        for label, rxs, now in compiled:
            for rx in rxs:
                m = rx.search(txt)
                if m:
                    low = txt.lower()
                    if any(a in low for a in ALLOW):
                        continue          # a documented correction, not a stale value
                    s = max(0, m.start() - 55)
                    hits.append((where, label, now, txt[s:m.end() + 45].strip()))
                    break

    name = os.path.basename(path)
    if not hits:
        print(f"  [ok  ] {name}: no superseded values")
        return 0
    print(f"  [FAIL] {name}: {len(hits)} superseded value(s)")
    for where, label, now, ctx in hits:
        print(f"           {where}: {label} — should read {now}")
        print(f"             ...{ctx}...")
    return len(hits)


DEFAULTS = [
    "/mnt/c/Users/user/Desktop/paperinfo-slugs_hydrates/paper5.docx",
    "/mnt/c/Users/user/Desktop/paperinfo-slugs_hydrates/paper5_typeset.docx",
    "/mnt/c/Users/user/Desktop/paperinfo-slugs_hydrates/paperinfo2-slugs.docx",
    "/mnt/c/Users/user/Desktop/personal projects/slugs_hydrates/slugs1-thesis.docx",
    "/mnt/c/Users/user/Desktop/slides3.pptx",
]


def main(argv):
    paths = argv or [p for p in DEFAULTS if os.path.exists(p)]
    if not paths:
        print("no documents found to check")
        return 2
    print(f"=== superseded-value check over {len(paths)} document(s) ===")
    print(f"    {len(RETIRED)} retired quantities, all known spellings, "
          f"table cells searched individually\n")
    total = sum(check(p) for p in paths)

    #  the exported PDFs, which are what actually get submitted
    stale = check_pdf_freshness(sibling_docs(paths))
    print()
    if stale:
        print(f"  [FAIL] {len(stale)} PDF(s) older than the document they were exported from:")
        for pdf, hours in stale:
            print(f"           {os.path.basename(pdf)} — {hours:.1f} h behind its .docx; re-export it")
        total += len(stale)
    else:
        print("  [ok  ] every exported PDF is at least as new as its source document")

    print()
    if total:
        print(f"{total} superseded value(s) — the documents disagree with the current outputs")
        return 1
    print("Every document is consistent with the current solver outputs.")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
