#!/usr/bin/env python3
# =============================================================================
#  audit_deck.py — slide-by-slide report of everything that makes a deck
#  unreadable, including the one thing geometry checks cannot see.
# -----------------------------------------------------------------------------
#  check_slides.py answers "is the layout sound?" — empty frames, overlaps,
#  truncation, stale figures. It passed a deck whose figures were still illegible,
#  because the fault was not in the layout: it was that a figure drawn for a 3.5 in
#  journal column had been placed in a 3 in frame, shrinking its 10 pt axis labels
#  to about 4 pt on the wall. Nothing about the shape geometry reveals that.
#
#  The missing quantity is the EFFECTIVE type size — how large the text inside a
#  figure actually renders once the figure is scaled into its frame:
#
#      effective_pt = base_pt x (displayed_width / natural_width)
#      natural_width = pixel_width / dpi
#
#  base_pt is 10 for a figure rendered for print and 18 for one rendered with
#  SHCT_FIG_FONTSCALE=1.8, so the set a figure came from has to be identified, and
#  it is, by hashing the embedded bytes against the figure directories on disk.
#
#  A projected slide needs roughly 12 pt at slide scale to be read from the back of
#  a room; below about 8 pt is not readable at all. Both thresholds are reported so
#  the marginal cases can be judged rather than merely counted.
#
#      python3 audit_deck.py [deck.pptx]
#
#  Exit status is 1 if anything is flagged.
# =============================================================================
import hashlib
import io as _io
import os
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
CASE = os.path.abspath(os.path.join(HERE, ".."))

READABLE_PT = 12.0        # comfortable from the back of a room
MARGINAL_PT = 8.0         # below this, effectively unreadable
MIN_FIG_SQIN = 6.0        # a figure smaller than this is a thumbnail
#  every slide-legible variant set, at whatever size scale it was rendered
SCALED_DIRS = tuple(sorted(
    d for d in os.listdir(CASE)
    if d.startswith("outputs_slides") and os.path.isdir(os.path.join(CASE, d))))
PRINT_DIRS = ("outputs_steady", "outputs_shutin", "outputs_mitigated")

#  A GIF carries no resolution metadata, so its natural size cannot be read back
#  from the file. Assuming the PIL default of 100 dpi understates every animation's
#  natural width by a third and made them look worse in this audit than they are.
#  make_animations.py renders at SHCT_ANIM_DPI, so use that.
_ANIM_DPI_DEFAULT = float(os.environ.get("SHCT_ANIM_DPI", "150"))


def anim_dpi(path):
    """The dpi an animation was rendered at, from the sidecar beside it."""
    import json
    d = os.path.dirname(path) or "."
    side = os.path.join(d, "anim_dpi.json")
    try:
        with open(side) as fh:
            return float(json.load(fh)["dpi"])
    except Exception:
        return _ANIM_DPI_DEFAULT


ANIM_DPI = _ANIM_DPI_DEFAULT


def box(sh):
    return (Emu(sh.left).inches, Emu(sh.top).inches,
            Emu(sh.width).inches, Emu(sh.height).inches)


def overlap_frac(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    x0, x1 = max(ax, bx), min(ax + aw, bx + bw)
    y0, y1 = max(ay, by), min(ay + ah, by + bh)
    if x1 <= x0 or y1 <= y0:
        return 0.0
    return (x1 - x0) * (y1 - y0) / max(min(aw * ah, bw * bh), 1e-9)


def figure_index():
    """sha256 -> (basename, base_pt) for every figure the deck could carry."""
    idx = {}
    for dirs, base_pt in ((PRINT_DIRS, 10.0), (SCALED_DIRS, 18.0)):
        for d in dirs:
            p = os.path.join(CASE, d)
            if not os.path.isdir(p):
                continue
            for f in os.listdir(p):
                if f.endswith((".png", ".gif")):
                    try:
                        h = hashlib.sha256(open(os.path.join(p, f), "rb").read()).hexdigest()
                    except OSError:
                        continue
                    #  the path is kept so a GIF's render dpi can be looked up
                    #  from the sidecar beside it
                    idx[h] = (f, base_pt, os.path.join(p, f))
    return idx


def natural_inches(blob, gif_dpi=None):
    """Width in inches at the dpi the figure was rendered with.

    A GIF stores no resolution, so its dpi is supplied by the caller from the
    sidecar make_animations.py writes beside it.
    """
    gif_dpi = _ANIM_DPI_DEFAULT if gif_dpi is None else gif_dpi
    im = Image.open(_io.BytesIO(blob))
    px = im.size[0]
    dpi = gif_dpi if (im.format or "").upper() == "GIF" else 100.0
    try:
        d = im.info.get("dpi")
        if d and d[0] > 1:
            dpi = float(d[0])
    except Exception:
        pass
    return px / dpi, px, dpi


def audit(path):
    prs = Presentation(path)
    SW, SH = Emu(prs.slide_width).inches, Emu(prs.slide_height).inches
    idx = figure_index()
    n_issue = 0
    worst = []

    print(f"=== {os.path.basename(path)} — {len(prs.slides._sldIdLst)} slides, "
          f"{SW:.2f} x {SH:.2f} in ===\n")

    for i, slide in enumerate(prs.slides, 1):
        pics = [s for s in slide.shapes if s.__class__.__name__ == "Picture"]
        texts = [s for s in slide.shapes
                 if s.has_text_frame and s.text_frame.text.strip()]
        issues = []

        for p in pics:
            l, t, w, h = box(p)
            try:
                blob = p.image.blob
            except Exception:
                continue
            name, base_pt, src = idx.get(hashlib.sha256(blob).hexdigest(),
                                         ("(unknown)", 10.0, None))
            try:
                nat_w, px, dpi = natural_inches(blob,
                                                anim_dpi(src) if src else None)
            except Exception:
                continue
            eff = base_pt * (w / nat_w) if nat_w else 0.0
            worst.append((eff, i, name))

            if w * h < MIN_FIG_SQIN:
                issues.append(f"{name}: only {w:.2f}x{h:.2f} in ({w*h:.1f} sq in)")
            if eff < MARGINAL_PT:
                issues.append(f"{name}: text renders at ~{eff:.1f} pt on the slide "
                              f"— not readable (figure is {nat_w:.1f} in natural, "
                              f"shown at {w:.2f} in)")
            elif eff < READABLE_PT:
                issues.append(f"{name}: text renders at ~{eff:.1f} pt — marginal from "
                              f"the back of a room")
            if l < -0.02 or t < -0.02 or l + w > SW + 0.02 or t + h > SH + 0.02:
                issues.append(f"{name}: runs off the slide "
                              f"(x {l:.2f}..{l+w:.2f}, y {t:.2f}..{t+h:.2f})")

        for a in range(len(pics)):
            for b in range(a + 1, len(pics)):
                f = overlap_frac(box(pics[a]), box(pics[b]))
                if f > 0.01:
                    issues.append(f"two figures overlap by {f*100:.0f} %")

        for tx in texts:
            for p in pics:
                f = overlap_frac(box(tx), box(p))
                if f >= 0.10:
                    issues.append(f"text {tx.text_frame.text.strip()[:32]!r} covers "
                                  f"{f*100:.0f} % of a figure")

        if issues:
            print(f"  slide {i:2d} ({len(pics)} figure(s))")
            for m in issues:
                print(f"        - {m}")
                n_issue += 1

    worst.sort()
    print(f"\n  least readable figures (effective type size on the slide):")
    for eff, i, name in worst[:8]:
        tag = "unreadable" if eff < MARGINAL_PT else ("marginal" if eff < READABLE_PT else "ok")
        print(f"     slide {i:2d}  {eff:5.1f} pt  {tag:11s} {name}")
    ok = sum(1 for e, _, _ in worst if e >= READABLE_PT)
    print(f"\n  {ok}/{len(worst)} figures render at {READABLE_PT:.0f} pt or better")
    print(f"  {n_issue} issue(s)")
    return 1 if n_issue else 0


def main(argv):
    path = argv[0] if argv else "/mnt/c/Users/user/Desktop/slides3.pptx"
    if not os.path.exists(path):
        print(f"not found: {path}")
        return 2
    return audit(path)


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
